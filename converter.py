# """import os
# import PyPDF2
# from PIL import Image
# from reportlab.pdfgen import canvas
# from reportlab.lib.pagesizes import letter
# from docx import Document
# from pptx import Presentation
# import tempfile

# def pdf_to_text(input_path, output_path):
#     """Convert PDF to text file"""
#     try:
#         with open(input_path, 'rb') as pdf_file:
#             reader = PyPDF2.PdfReader(pdf_file)
#             text = ""
#             for page in reader.pages:
#                 text += page.extract_text() + "\n"
            
#             with open(output_path, 'w', encoding='utf-8') as text_file:
#                 text_file.write(text)
#         return True
#     except Exception as e:
#         print(f"Error in pdf_to_text: {e}")
#         return False

# def image_to_pdf(input_path, output_path):
#     """Convert image to PDF"""
#     try:
#         image = Image.open(input_path)
#         if image.mode != 'RGB':
#             image = image.convert('RGB')
        
#         image.save(output_path, 'PDF')
#         return True
#     except Exception as e:
#         print(f"Error in image_to_pdf: {e}")
#         return False

# def text_to_pdf(input_path, output_path):
#     """Convert text file to PDF"""
#     try:
#         with open(input_path, 'r', encoding='utf-8') as text_file:
#             text_content = text_file.read()
        
#         c = canvas.Canvas(output_path, pagesize=letter)
#         text = c.beginText(40, 750)
#         text.setFont("Helvetica", 12)
        
#         lines = []
#         for line in text_content.split('\n'):
#             while len(line) > 90:
#                 lines.append(line[:90])
#                 line = line[90:]
#             lines.append(line)
        
#         for line in lines:
#             if text.getY() < 40:
#                 c.drawText(text)
#                 c.showPage()
#                 text = c.beginText(40, 750)
#                 text.setFont("Helvetica", 12)
#             text.textLine(line)
        
#         c.drawText(text)
#         c.save()
#         return True
#     except Exception as e:
#         print(f"Error in text_to_pdf: {e}")
#         return False

# def pdf_to_docx(input_path, output_path):
#     """Convert PDF to DOCX (simple text extraction approach)"""
#     try:
#         # First extract text from PDF
#         with open(input_path, 'rb') as pdf_file:
#             reader = PyPDF2.PdfReader(pdf_file)
#             text_content = ""
#             for page in reader.pages:
#                 text_content += page.extract_text() + "\n\n"
        
#         # Create DOCX from extracted text
#         doc = Document()
#         for paragraph in text_content.split('\n\n'):
#             if paragraph.strip():
#                 doc.add_paragraph(paragraph)
        
#         doc.save(output_path)
#         return True
#     except Exception as e:
#         print(f"Error in pdf_to_docx: {e}")
#         return False

# def docx_to_pdf(input_path, output_path):
#     """Convert DOCX to PDF (using text extraction approach)"""
#     try:
#         doc = Document(input_path)
#         text_content = ""
#         for paragraph in doc.paragraphs:
#             text_content += paragraph.text + "\n"
        
#         c = canvas.Canvas(output_path, pagesize=letter)
#         text = c.beginText(40, 750)
#         text.setFont("Helvetica", 12)
        
#         lines = []
#         for line in text_content.split('\n'):
#             while len(line) > 90:
#                 lines.append(line[:90])
#                 line = line[90:]
#             lines.append(line)
        
#         for line in lines:
#             if text.getY() < 40:
#                 c.drawText(text)
#                 c.showPage()
#                 text = c.beginText(40, 750)
#                 text.setFont("Helvetica", 12)
#             text.textLine(line)
        
#         c.drawText(text)
#         c.save()
#         return True
#     except Exception as e:
#         print(f"Error in docx_to_pdf: {e}")
#         return False

# def docx_to_text(input_path, output_path):
#     """Convert DOCX to text"""
#     try:
#         doc = Document(input_path)
#         text_content = ""
#         for paragraph in doc.paragraphs:
#             text_content += paragraph.text + "\n"
        
#         with open(output_path, 'w', encoding='utf-8') as text_file:
#             text_file.write(text_content)
#         return True
#     except Exception as e:
#         print(f"Error in docx_to_text: {e}")
#         return False

# def text_to_docx(input_path, output_path):
#     """Convert text to DOCX"""
#     try:
#         with open(input_path, 'r', encoding='utf-8') as text_file:
#             text_content = text_file.read()
        
#         doc = Document()
#         for line in text_content.split('\n'):
#             doc.add_paragraph(line)
        
#         doc.save(output_path)
#         return True
#     except Exception as e:
#         print(f"Error in text_to_docx: {e}")
#         return False

# def pptx_to_text(input_path, output_path):
#     """Convert PPTX to text"""
#     try:
#         prs = Presentation(input_path)
#         text_content = ""
        
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text_content += shape.text + "\n"
        
#         with open(output_path, 'w', encoding='utf-8') as text_file:
#             text_file.write(text_content)
#         return True
#     except Exception as e:
#         print(f"Error in pptx_to_text: {e}")
#         return False

# def text_to_pptx(input_path, output_path):
#     """Convert text to PPTX"""
#     try:
#         with open(input_path, 'r', encoding='utf-8') as text_file:
#             text_content = text_file.read()
        
#         prs = Presentation()
        
#         # Split text into slides (each paragraph becomes a slide)
#         paragraphs = text_content.split('\n\n')
#         for para in paragraphs:
#             if para.strip():
#                 slide = prs.slides.add_slide(prs.slide_layouts[1])
#                 title = slide.shapes.title
#                 content = slide.placeholders[1]
                
#                 title.text = "Slide"
#                 content.text = para
        
#         prs.save(output_path)
#         return True
#     except Exception as e:
#         print(f"Error in text_to_pptx: {e}")
#         return False"""

# ==========================
# IMPORT REQUIRED LIBRARIES
# ==========================
import os  # For file path operations
import PyPDF2  # For reading PDF files
from PIL import Image  # For image processing
from reportlab.pdfgen import canvas  # For generating PDF files
from reportlab.lib.pagesizes import letter  # Standard page size for PDFs
from docx import Document  # For working with Word DOCX files
from pptx import Presentation  # For reading/writing PowerPoint files
import tempfile  # For creating temporary files (not used here but useful)

# =====================
# PDF to TEXT FUNCTION
# =====================
def pdf_to_text(input_path, output_path):
    """Convert PDF to plain text"""
    try:
        with open(input_path, 'rb') as pdf_file:  # Open PDF file in binary read mode
            reader = PyPDF2.PdfReader(pdf_file)  # Create PDF reader object
            text = ""
            for page in reader.pages:  # Loop through each page of the PDF
                text += page.extract_text() + "\n"  # Extract and add text
        
        with open(output_path, 'w', encoding='utf-8') as text_file:  # Open output text file
            text_file.write(text)  # Write extracted text
        return True
    except Exception as e:
        print(f"Error in pdf_to_text: {e}")  # Print error message
        return False

# =====================
# IMAGE to PDF FUNCTION
# =====================
def image_to_pdf(input_path, output_path):
    """Convert image to PDF"""
    try:
        image = Image.open(input_path)  # Open the image
        if image.mode != 'RGB':  # Check if the image is not in RGB mode
            image = image.convert('RGB')  # Convert to RGB
        
        image.save(output_path, 'PDF')  # Save the image as a PDF
        return True
    except Exception as e:
        print(f"Error in image_to_pdf: {e}")
        return False

# =====================
# TEXT to PDF FUNCTION
# =====================
def text_to_pdf(input_path, output_path):
    """Convert plain text file to PDF"""
    try:
        with open(input_path, 'r', encoding='utf-8') as text_file:
            text_content = text_file.read()  # Read all text content

        c = canvas.Canvas(output_path, pagesize=letter)  # Create PDF canvas
        text = c.beginText(40, 750)  # Set initial text position (x=40, y=750)
        text.setFont("Helvetica", 12)  # Set font and size

        lines = []  # List to store wrapped lines
        for line in text_content.split('\n'):  # Process each line
            while len(line) > 90:  # If line too long, wrap at 90 chars
                lines.append(line[:90])
                line = line[90:]
            lines.append(line)  # Add remaining part

        for line in lines:
            if text.getY() < 40:  # If we reach the bottom of the page
                c.drawText(text)  # Draw current text block
                c.showPage()  # Add new page
                text = c.beginText(40, 750)  # Reset text position
                text.setFont("Helvetica", 12)
            text.textLine(line)  # Add line to text object

        c.drawText(text)  # Final draw
        c.save()  # Save the PDF
        return True
    except Exception as e:
        print(f"Error in text_to_pdf: {e}")
        return False

# =====================
# PDF to DOCX FUNCTION
# =====================
def pdf_to_docx(input_path, output_path):
    """Convert PDF to DOCX using extracted text"""
    try:
        with open(input_path, 'rb') as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            text_content = ""
            for page in reader.pages:
                text_content += page.extract_text() + "\n\n"  # Extract page text

        doc = Document()  # Create new Word document
        for paragraph in text_content.split('\n\n'):  # Split into paragraphs
            if paragraph.strip():  # Ignore empty ones
                doc.add_paragraph(paragraph)

        doc.save(output_path)  # Save DOCX
        return True
    except Exception as e:
        print(f"Error in pdf_to_docx: {e}")
        return False

# =====================
# DOCX to PDF FUNCTION
# =====================
def docx_to_pdf(input_path, output_path):
    """Convert DOCX to PDF by rendering text"""
    try:
        doc = Document(input_path)  # Load Word file
        text_content = ""
        for paragraph in doc.paragraphs:  # Loop through all paragraphs
            text_content += paragraph.text + "\n"

        c = canvas.Canvas(output_path, pagesize=letter)
        text = c.beginText(40, 750)
        text.setFont("Helvetica", 12)

        lines = []
        for line in text_content.split('\n'):
            while len(line) > 90:
                lines.append(line[:90])
                line = line[90:]
            lines.append(line)

        for line in lines:
            if text.getY() < 40:
                c.drawText(text)
                c.showPage()
                text = c.beginText(40, 750)
                text.setFont("Helvetica", 12)
            text.textLine(line)

        c.drawText(text)
        c.save()
        return True
    except Exception as e:
        print(f"Error in docx_to_pdf: {e}")
        return False

# =====================
# DOCX to TEXT FUNCTION
# =====================
def docx_to_text(input_path, output_path):
    """Convert DOCX to plain text"""
    try:
        doc = Document(input_path)  # Load Word file
        text_content = ""
        for paragraph in doc.paragraphs:  # Read all paragraphs
            text_content += paragraph.text + "\n"

        with open(output_path, 'w', encoding='utf-8') as text_file:
            text_file.write(text_content)  # Save as plain text
        return True
    except Exception as e:
        print(f"Error in docx_to_text: {e}")
        return False

# =====================
# TEXT to DOCX FUNCTION
# =====================
def text_to_docx(input_path, output_path):
    """Convert plain text file to DOCX"""
    try:
        with open(input_path, 'r', encoding='utf-8') as text_file:
            text_content = text_file.read()

        doc = Document()  # Create Word document
        for line in text_content.split('\n'):  # Each line becomes a paragraph
            doc.add_paragraph(line)

        doc.save(output_path)
        return True
    except Exception as e:
        print(f"Error in text_to_docx: {e}")
        return False

# =====================
# PPTX to TEXT FUNCTION
# =====================
def pptx_to_text(input_path, output_path):
    """Convert PPTX to plain text"""
    try:
        prs = Presentation(input_path)  # Load PowerPoint file
        text_content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):  # Check if shape contains text
                    text_content += shape.text + "\n"

        with open(output_path, 'w', encoding='utf-8') as text_file:
            text_file.write(text_content)
        return True
    except Exception as e:
        print(f"Error in pptx_to_text: {e}")
        return False

# =====================
# TEXT to PPTX FUNCTION
# =====================
def text_to_pptx(input_path, output_path):
    """Convert plain text file to PPTX slides"""
    try:
        with open(input_path, 'r', encoding='utf-8') as text_file:
            text_content = text_file.read()

        prs = Presentation()  # Create new PowerPoint presentation
        paragraphs = text_content.split('\n\n')  # Each paragraph becomes a slide

        for para in paragraphs:
            if para.strip():  # Ignore empty paragraphs
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content slide
                title = slide.shapes.title  # Slide title placeholder
                content = slide.placeholders[1]  # Content placeholder
                title.text = "Slide"  # Set title
                content.text = para  # Set content

        prs.save(output_path)  # Save PPTX file
        return True
    except Exception as e:
        print(f"Error in text_to_pptx: {e}")
        return False

# =====================
# SAMPLE USAGE SECTION
# =====================
if __name__ == "__main__":
    # Example usage (replace these paths with your actual file paths)
    pdf_to_text("sample.pdf", "output.txt")
    image_to_pdf("image.jpg", "image.pdf")
    text_to_pdf("notes.txt", "notes.pdf")
    pdf_to_docx("sample.pdf", "output.docx")
    docx_to_pdf("document.docx", "document.pdf")
    docx_to_text("document.docx", "document.txt")
    text_to_docx("notes.txt", "notes.docx")
   
